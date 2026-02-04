from __future__ import annotations

import json
import re
import time
import zipfile
from dataclasses import dataclass
from pathlib import Path

from .init_project import init_project


@dataclass(frozen=True)
class ImportedFile:
    project_dir: Path
    in_path: Path
    md_relpath: str
    copied_source_relpath: str


def import_file(
    in_path: Path,
    *,
    slug: str | None = None,
    title: str | None = None,
    force: bool = False,
    extract_images: bool = True,
) -> ImportedFile:
    in_path = in_path.expanduser().resolve()
    if not in_path.exists():
        raise FileNotFoundError(f"input not found: {in_path}")
    if in_path.is_dir():
        raise IsADirectoryError(f"input must be a file: {in_path}")

    suffix = in_path.suffix.lower()
    if suffix not in {".md", ".txt", ".docx", ".pptx"}:
        raise ValueError(f"unsupported file type: {suffix} (supported: .md .txt .docx .pptx)")

    slug = slug or _suggest_slug_from_filename(in_path.name)
    project_dir = init_project(slug).resolve()

    md_path = project_dir / "source" / "article.md"
    if md_path.exists() and not force:
        raise FileExistsError(f"already exists: {md_path} (use --force to overwrite)")

    copied = project_dir / "source" / f"input{suffix}"
    copied.write_bytes(in_path.read_bytes())

    assets_images_src = project_dir / "assets" / "images" / "source"
    if extract_images and suffix in {".docx", ".pptx"}:
        _extract_office_images(in_path, assets_images_src)

    if suffix == ".md":
        md = in_path.read_text(encoding="utf-8", errors="replace")
        md_path.write_text(md if md.endswith("\n") else md + "\n", encoding="utf-8")
    elif suffix == ".txt":
        text = in_path.read_text(encoding="utf-8", errors="replace").strip()
        md = _text_to_markdown(text, title=title or slug, images_dir=assets_images_src if assets_images_src.exists() else None)
        md_path.write_text(md, encoding="utf-8")
    elif suffix == ".docx":
        md = _docx_to_markdown(in_path, title=title or slug, images_dir=assets_images_src if assets_images_src.exists() else None)
        md_path.write_text(md, encoding="utf-8")
    elif suffix == ".pptx":
        md = _pptx_to_markdown(in_path, title=title or slug, images_dir=assets_images_src if assets_images_src.exists() else None)
        md_path.write_text(md, encoding="utf-8")

    _maybe_update_project_title(project_dir, title=title or slug, slug=slug)

    return ImportedFile(
        project_dir=project_dir,
        in_path=in_path,
        md_relpath="source/article.md",
        copied_source_relpath=f"source/input{suffix}",
    )


def _suggest_slug_from_filename(name: str) -> str:
    stem = Path(name).stem.strip().lower()
    stem = re.sub(r"[^0-9a-zA-Zぁ-んァ-ヶ一-龠ー_-]+", "_", stem).strip("_")
    return stem or "project"


def _maybe_update_project_title(project_dir: Path, *, title: str, slug: str) -> None:
    project_json = project_dir / "project.json"
    if not project_json.exists():
        return
    try:
        data = json.loads(project_json.read_text(encoding="utf-8"))
    except Exception:
        return
    if not isinstance(data, dict):
        return
    if data.get("title") not in (None, "", slug):
        return
    data["title"] = title
    project_json.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _frontmatter(*, source: str, title: str) -> str:
    return (
        "---\n"
        f"source_file: {source}\n"
        f"imported_at: {time.strftime('%Y-%m-%dT%H:%M:%S%z')}\n"
        "---\n\n"
        f"# {title}\n\n"
    )


def _text_to_markdown(text: str, *, title: str, images_dir: Path | None) -> str:
    out = [_frontmatter(source="text", title=title)]
    if images_dir and images_dir.exists():
        out.append(_images_section(images_dir))
    out.append(text.strip() + "\n")
    return "".join(out)


def _docx_to_markdown(docx_path: Path, *, title: str, images_dir: Path | None) -> str:
    _require_python_docx()
    from docx import Document  # type: ignore[import-not-found]

    doc = Document(str(docx_path))
    out: list[str] = [_frontmatter(source=str(docx_path.name), title=title)]
    if images_dir and images_dir.exists():
        out.append(_images_section(images_dir))

    for p in doc.paragraphs:
        txt = (p.text or "").strip()
        if not txt:
            continue
        style = (p.style.name or "").lower() if getattr(p, "style", None) is not None else ""
        level = _docx_heading_level(style)
        if level:
            out.append(f"{'#' * level} {txt}\n\n")
        else:
            out.append(txt + "\n\n")
    return "".join(out)


def _docx_heading_level(style_name_lower: str) -> int | None:
    # Common: "Heading 1", "Heading 2", ... ; also allow Japanese names heuristically.
    m = re.search(r"heading\\s*(\\d+)", style_name_lower)
    if m:
        lvl = int(m.group(1))
        return max(1, min(6, lvl))
    if "見出し" in style_name_lower:
        # e.g. 見出し 1
        m2 = re.search(r"(\\d+)", style_name_lower)
        if m2:
            lvl = int(m2.group(1))
            return max(1, min(6, lvl))
        return 2
    return None


def _pptx_to_markdown(pptx_path: Path, *, title: str, images_dir: Path | None) -> str:
    _require_python_pptx()
    from pptx import Presentation  # type: ignore[import-not-found]

    pres = Presentation(str(pptx_path))
    out: list[str] = [_frontmatter(source=str(pptx_path.name), title=title)]
    if images_dir and images_dir.exists():
        out.append(_images_section(images_dir))

    for idx, slide in enumerate(pres.slides, start=1):
        slide_title = _pptx_slide_title(slide) or f"Slide {idx}"
        out.append(f"## {slide_title}\n\n")
        bullets = _pptx_slide_bullets(slide)
        for b in bullets:
            out.append(f"- {b}\n")
        out.append("\n")
    return "".join(out)


def _pptx_slide_title(slide) -> str | None:
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            t = slide.shapes.title.text_frame.text.strip()
            return t or None
    except Exception:
        pass
    return None


def _pptx_slide_bullets(slide) -> list[str]:
    out: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for p in tf.paragraphs:
            txt = (p.text or "").strip()
            if not txt:
                continue
            # Skip title duplication.
            if shape == getattr(slide.shapes, "title", None):
                continue
            out.append(txt)
    # de-dupe while preserving order
    seen: set[str] = set()
    deduped: list[str] = []
    for x in out:
        if x in seen:
            continue
        seen.add(x)
        deduped.append(x)
    return deduped


def _images_section(images_dir: Path) -> str:
    # Put extracted images at top so they can be picked up later if needed.
    rels = sorted([p for p in images_dir.iterdir() if p.is_file()])
    if not rels:
        return ""
    lines = ["## 添付画像（抽出）\n\n"]
    for p in rels:
        lines.append(f"![]({p.as_posix()})\n\n")
    return "".join(lines)


def _extract_office_images(in_path: Path, out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(in_path) as z:
        names = z.namelist()
        if in_path.suffix.lower() == ".docx":
            media_prefix = "word/media/"
        else:
            media_prefix = "ppt/media/"

        extracted = 0
        for name in names:
            if not name.startswith(media_prefix):
                continue
            filename = Path(name).name
            if not filename:
                continue
            data = z.read(name)
            (out_dir / filename).write_bytes(data)
            extracted += 1
    if extracted == 0:
        # Keep directory empty if nothing was found.
        try:
            out_dir.rmdir()
        except OSError:
            pass


def _require_python_docx() -> None:
    try:
        import docx  # noqa: F401
    except Exception as e:
        raise RuntimeError(
            "docx の取り込みには python-docx が必要です。\n"
            "インストール: python -m pip install -r requirements.txt"
        ) from e


def _require_python_pptx() -> None:
    try:
        import pptx  # noqa: F401
    except Exception as e:
        raise RuntimeError(
            "pptx の取り込みには python-pptx が必要です。\n"
            "インストール: python -m pip install -r requirements.txt"
        ) from e
